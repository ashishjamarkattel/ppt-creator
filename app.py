import os
import openai
from pptx import Presentation
from pptx.shapes.placeholder import PicturePlaceholder
import random 
import re
from serpapi import GoogleSearch
import requests
import streamlit as st
import urllib.request
from requests import get

openai.api_key = ""
SERP_API = ""


Prompt = """Write a presentation/powerpoint about the user's topic. You only answer with the presentation. Follow the structure of the example.
Notice
-You do all the presentation text for the user.
-You write the texts no longer than 250 characters!
-You make very short titles!
-You make the presentation easy to understand.
-The presentation has a table of contents.
-The presentation has a summary.
-At least 6 slides.

Example! - Stick to this formatting exactly!
#Title: TITLE OF THE PRESENTATION

#Slide: 1
#Header: table of contents
#Content: 1. CONTENT OF THIS POWERPOINT
2. CONTENTS OF THIS POWERPOINT
3. CONTENT OF THIS POWERPOINT
...

#Slide: 2
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE

#Slide: 3
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE

#Slide: 4
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE

#Slide: 5
#Headers: summary
#Content: CONTENT OF THE SUMMARY

#Slide: END"""





## create ppt content via openai
def create_ppt_text(topic):
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": (Prompt)},
            {"role": "user", "content": ("The user wants a presentation about " + topic)}
        ],
        temperature=0.5,
    )
    print(response)
    return response['choices'][0]['message']['content']



## Function to download image vai serpapi
## pip install google-serp-api
def download_image(query):
    params = {
        "q": query,
        "engine": "google_images",
        "ijn": "0",
        "api_key": SERP_API
        }

    search = GoogleSearch(params)
    results = search.get_dict()
    images_results = results["images_results"][0]
    if images_results["original"]:
        response = requests.get(images_results["original"])
        with open(f"images/{query}.jpg", "wb") as f:
            f.write(response.content)



## write the generated content in ppt via pptx
## pip install python-pptx
def create_ppt(text_file, design_number, ppt_name, user_name):
    prs = Presentation(f"Designs/Design-{design_number}.pptx")
    slide_count = 0
    header = ""
    content = ""
    last_slide_layout_index = -1
    firsttime = True
    with open(text_file, 'r', encoding='utf-8') as f:
        for line_num, line in enumerate(f):
            if line.startswith('#Title:'):
                header = line.replace('#Title:', '').strip()
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title = slide.shapes.title
                body_shape = slide.shapes.placeholders[1]
                tf = body_shape.text_frame
                tf.text = f"Name of the maker: {user_name}"
              
                title.text = header

                body_shape = slide.shapes.placeholders[1]
                continue
            elif line.startswith('#Slide:'):
                if slide_count > 0:
                    slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
                    for i in slide.placeholders:
                        print(f"place holder name {i.name}, slide number: {slide_count}")

                        if "Text Placeholder" in i.name:
                            i.text = "Defination of something"
                            # i.insert_picture(r"images\cow.jpg")

                        if isinstance(i, PicturePlaceholder):
                            image_path = f"images\{header}.jpg"
                            # create_image(header=header)
                            download_image(header)
                            print(i.left, i.height.inches, i.width.inches)
                            if os.path.exists(image_path):
                                print(image_path)
                                print("image found")
                                try:
                                    print(i.left, i.height, i.width)
                                    picture = i.insert_picture(image_path)
                                    picture.crop_top = 0
                                    picture.crop_left = 0
                                    picture.crop_bottom = 0
                                    picture.crop_right = 0

                                except Exception as e:
                                    print(e)
                                    pass
                        
                
                    title = slide.shapes.title
                    title.text = header
                    body_shape = slide.shapes.placeholders[slide_placeholder_index]
                    tf = body_shape.text_frame
                    tf.text = content
     
                content = "" 
                slide_count += 1
                slide_layout_index = last_slide_layout_index
                layout_indices = [1,  8] 
                while slide_layout_index == last_slide_layout_index:
                    if firsttime == True:
                        slide_layout_index = 1
                        slide_placeholder_index = 1
                        firsttime = False
                        break
                    slide_layout_index = random.choice(layout_indices) # Select random slide index
                    if slide_layout_index == 8:
                        slide_placeholder_index = 2
                    else:
                        slide_placeholder_index = 1
                last_slide_layout_index = slide_layout_index
                continue

            elif line.startswith('#Header:'):
                header = line.replace('#Header:', '').strip()
                continue

            elif line.startswith('#Content:'):
                content = line.replace('#Content:', '').strip()
                next_line = f.readline().strip()
                while next_line and not next_line.startswith('#'):
                    content += '\n' + next_line
                    next_line = f.readline().strip()
                continue

    prs.save(f'GeneratedPresentations/{ppt_name}.pptx')
    file_path = f"GeneratedPresentations/{ppt_name}.pptx"

    return file_path


## wrapper
def main(topic, user_name):

    input_string = topic
    input_string = re.sub(r'[^\w\s.\-\(\)]', '', input_string)
    input_string = input_string.replace("\n", "")
    design_pptx = [i for i in range(1, 10)]
    number = random.choice(design_pptx)

    if not os.path.exists(f'Cache/{input_string}.txt'):
        with open(f'Cache/{input_string}.txt', 'w', encoding='utf-8') as f:
            f.write(create_ppt_text(input_string))

    generated_ppt_path = create_ppt(f'Cache/{topic}.txt', number, input_string, user_name)
    return generated_ppt_path


if __name__ == "__main__":

    st.set_page_config(page_title= "Powerpoint Generator", page_icon= ":bird:")
    st.header("Generate Presentation :bird:")

    ppt_topic = st.text_input("PPTNAME")
    user_name = st.text_input("USERNAME")
    button = st.button("Generate PPT")

    if button:
        # Process the input
        with st.spinner("Processing input..."):
            generated_ppt_path = main(topic=ppt_topic, user_name=user_name)
            with open(generated_ppt_path, "rb") as f:
                st.download_button('Download', f, file_name=f"{ppt_topic}.pptx")
    

